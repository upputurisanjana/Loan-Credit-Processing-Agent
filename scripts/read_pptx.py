import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation

prs = Presentation("Credit_Decisioning_Agent.pptx")
print(f"Total slides: {len(prs.slides)}")
print()
for i, slide in enumerate(prs.slides, 1):
    layout = slide.slide_layout.name if slide.slide_layout else "?"
    print(f"=== SLIDE {i} ({layout}) ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            name = shape.name
            for j, para in enumerate(shape.text_frame.paragraphs):
                t = para.text.strip()
                if t:
                    level = para.level
                    indent = "  " * level
                    print(f"  [{name}] {indent}{t}")
    print()
